"""Document parsers for different file formats"""

import pandas as pd
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
import json

# Document processing imports
try:
    import openpyxl
    from docx import Document
    import PyPDF2
    from pptx import Presentation
    import fitz  # PyMuPDF
    import tabula
    import camelot
    # OCR imports
    import pytesseract
    import easyocr
    from pdf2image import convert_from_path
    from PIL import Image
except ImportError as e:
    logging.warning(f"Some document processing libraries not available: {e}")

from .file_detector import FileFormat, FileDetector

logger = logging.getLogger(__name__)


class DocumentParseError(Exception):
    """Custom exception for document parsing errors"""
    pass


class DocumentParser:
    """Base class for document parsers"""
    
    def parse(self, file_path: Path) -> Dict[str, Any]:
        """Parse document and return structured data"""
        raise NotImplementedError
    
    def to_dataframe(self, parsed_data: Dict[str, Any]) -> pd.DataFrame:
        """Convert parsed data to DataFrame for processing"""
        raise NotImplementedError


class CSVParser(DocumentParser):
    """Parser for CSV files"""
    
    def parse(self, file_path: Path) -> Dict[str, Any]:
        """Parse CSV file"""
        try:
            df = pd.read_csv(file_path)
            return {
                "type": "tabular",
                "data": df,
                "metadata": {
                    "rows": len(df),
                    "columns": len(df.columns),
                    "column_names": list(df.columns)
                }
            }
        except Exception as e:
            raise DocumentParseError(f"Error parsing CSV: {e}")
    
    def to_dataframe(self, parsed_data: Dict[str, Any]) -> pd.DataFrame:
        return parsed_data["data"]


class ExcelParser(DocumentParser):
    """Parser for Excel files (XLSX, XLS)"""
    
    def parse(self, file_path: Path) -> Dict[str, Any]:
        """Parse Excel file"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            sheets_data = {}
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                sheets_data[sheet_name] = df
            
            # If only one sheet, use it directly
            if len(sheets_data) == 1:
                main_df = list(sheets_data.values())[0]
            else:
                # Combine all sheets
                main_df = pd.concat(sheets_data.values(), ignore_index=True)
            
            return {
                "type": "tabular",
                "data": main_df,
                "sheets": sheets_data,
                "metadata": {
                    "sheet_names": excel_file.sheet_names,
                    "total_rows": len(main_df),
                    "total_columns": len(main_df.columns) if not main_df.empty else 0,
                    "column_names": list(main_df.columns) if not main_df.empty else []
                }
            }
        except Exception as e:
            raise DocumentParseError(f"Error parsing Excel: {e}")
    
    def to_dataframe(self, parsed_data: Dict[str, Any]) -> pd.DataFrame:
        return parsed_data["data"]


class PDFParser(DocumentParser):
    """Parser for PDF files"""
    
    def parse(self, file_path: Path) -> Dict[str, Any]:
        """Parse PDF file with OCR support"""
        try:
            # Try to extract tables first using camelot
            tables_data = self._extract_tables_camelot(file_path)
            
            # If no tables found, try tabula
            if not tables_data:
                tables_data = self._extract_tables_tabula(file_path)
            
            # Extract text content using PyMuPDF first
            text_content = self._extract_text_pymupdf(file_path)
            
            # Always try OCR for better text extraction and to save OCR output
            used_ocr = False
            logger.info("Attempting OCR extraction to save OCR output...")
            print("🔍 Attempting OCR extraction...")
            ocr_text = self._extract_text_ocr(file_path)
            
            # Always prefer OCR text for PDFs to ensure LLM gets the best quality text
            if len(ocr_text.strip()) > 100:  # OCR succeeded and has meaningful content
                text_content = ocr_text
                used_ocr = True
                print("✅ Using OCR text extraction for LLM processing")
                logger.info("Using OCR text extraction for LLM processing")
            elif len(text_content.strip()) > 100:  # Fallback to regular extraction
                print("⚠️ Using regular text extraction (OCR had insufficient content)")
                logger.info("Using regular text extraction (OCR had insufficient content)")
            else:
                print("❌ Both OCR and regular extraction failed to get sufficient text")
                logger.warning("Both OCR and regular extraction failed to get sufficient text")
            
            # If we have tables, use them; otherwise convert text to structured format
            if tables_data:
                main_df = pd.concat(tables_data, ignore_index=True)
                data_type = "tabular"
            else:
                # Convert text to simple structure
                lines = [line.strip() for line in text_content.split('\n') if line.strip()]
                main_df = pd.DataFrame({'content': lines})
                data_type = "text"
            
            return {
                "type": data_type,
                "data": main_df,
                "raw_text": text_content,
                "tables_found": len(tables_data) if tables_data else 0,
                "metadata": {
                    "pages": len(text_content.split('\n\n')),  # Rough page count
                    "rows": len(main_df),
                    "columns": len(main_df.columns),
                    "has_tables": bool(tables_data),
                    "used_ocr": used_ocr
                }
            }
        except Exception as e:
            raise DocumentParseError(f"Error parsing PDF: {e}")
    
    def _extract_tables_camelot(self, file_path: Path) -> List[pd.DataFrame]:
        """Extract tables using camelot"""
        try:
            tables = camelot.read_pdf(str(file_path), pages='all')
            return [table.df for table in tables if not table.df.empty]
        except Exception as e:
            logger.debug(f"Camelot table extraction failed: {e}")
            return []
    
    def _extract_tables_tabula(self, file_path: Path) -> List[pd.DataFrame]:
        """Extract tables using tabula"""
        try:
            tables = tabula.read_pdf(str(file_path), pages='all', multiple_tables=True)
            return [df for df in tables if not df.empty]
        except Exception as e:
            logger.debug(f"Tabula table extraction failed: {e}")
            return []
    
    def _extract_text_pymupdf(self, file_path: Path) -> str:
        """Extract text using PyMuPDF"""
        try:
            doc = fitz.open(str(file_path))
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            logger.debug(f"PyMuPDF text extraction failed: {e}")
            return ""
    
    def _extract_text_ocr(self, file_path: Path) -> str:
        """Enhanced OCR text extraction with improved quality and preprocessing"""
        try:
            # Convert PDF to images with optimized settings
            print(f"🔍 Converting PDF to images for OCR processing...")
            images = convert_from_path(
                str(file_path), 
                dpi=300,  # Higher DPI for better quality
                first_page=1, 
                last_page=10,  # Process more pages
                grayscale=True,  # Grayscale often works better for OCR
                thread_count=2   # Use multiple threads
            )
            
            ocr_text = ""
            
            # Try EasyOCR first with enhanced settings
            try:
                # Support multiple languages for better detection
                reader = easyocr.Reader(['en', 'it', 'fr', 'de', 'es'], gpu=False)
                
                for i, image in enumerate(images):
                    logger.info(f"Processing page {i+1}/{len(images)} with EasyOCR...")
                    print(f"🔍 OCR processing page {i+1}/{len(images)} with EasyOCR...")
                    
                    # Preprocess image for better OCR results
                    image = self._preprocess_image_for_ocr(image)
                    
                    # Convert PIL image to numpy array for EasyOCR
                    import numpy as np
                    img_array = np.array(image)
                    
                    # EasyOCR with enhanced parameters
                    results = reader.readtext(
                        img_array,
                        detail=1,  # Get detailed results
                        paragraph=True,  # Group text into paragraphs
                        width_ths=0.7,  # Text width threshold
                        height_ths=0.7,  # Text height threshold
                        min_size=20,  # Minimum text size
                        text_threshold=0.7,  # Text confidence threshold
                        low_text=0.4,  # Low text threshold
                        link_threshold=0.4,  # Link threshold
                        canvas_size=2560,  # Canvas size for better processing
                        mag_ratio=1.5  # Magnification ratio
                    )
                    
                    # Process results with better text consolidation
                    if results:
                        page_text = self._consolidate_ocr_results(results)
                        ocr_text += f"\n--- Page {i+1} ---\n{page_text}\n"
                    
                logger.info(f"EasyOCR extracted {len(ocr_text)} characters")
                print(f"✅ EasyOCR extracted {len(ocr_text)} characters")
                
            except Exception as e:
                logger.warning(f"EasyOCR failed: {e}, trying Tesseract...")
                print(f"⚠️ EasyOCR failed: {e}, trying Tesseract...")
                
                # Fallback to Tesseract with enhanced configuration
                for i, image in enumerate(images):
                    logger.info(f"Processing page {i+1}/{len(images)} with Tesseract...")
                    print(f"🔍 OCR processing page {i+1}/{len(images)} with Tesseract...")
                    
                    # Preprocess image for better OCR results
                    image = self._preprocess_image_for_ocr(image)
                    
                    # Enhanced Tesseract configuration (fixed quotation issue)
                    custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=0123456789abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ.,;:!?()[]{}«»"àáâãäåæçèéêëìíîïñòóôõöøùúûüýÿ'
                    page_text = pytesseract.image_to_string(image, lang='eng+ita+fra+deu+spa', config=custom_config)
                    ocr_text += f"\n--- Page {i+1} ---\n{page_text}\n"
                    
                logger.info(f"Tesseract extracted {len(ocr_text)} characters")
                print(f"✅ Tesseract extracted {len(ocr_text)} characters")
            
            # Enhanced cleaning and filtering
            if ocr_text.strip():
                cleaned_text = self._clean_ocr_text_enhanced(ocr_text)
                
                # Always save OCR output as TXT file (both raw and cleaned) in output/ocr folder
                output_ocr_dir = Path("output/ocr")
                output_ocr_dir.mkdir(parents=True, exist_ok=True)
                ocr_output_path = output_ocr_dir / f"{file_path.stem}.ocr.txt"
                with open(ocr_output_path, 'w', encoding='utf-8') as f:
                    f.write(f"OCR Output for: {file_path.name}\n")
                    f.write(f"Generated: {__import__('datetime').datetime.now()}\n")
                    f.write("=" * 50 + "\n\n")
                    f.write("=== RAW OCR OUTPUT ===\n")
                    f.write(ocr_text)
                    f.write("\n\n=== CLEANED OCR OUTPUT ===\n")
                    f.write(cleaned_text)
                
                print(f"💾 OCR output saved to: {ocr_output_path}")
                print(f"🧹 Cleaned text: {len(cleaned_text)} chars (was {len(ocr_text)})")
                logger.info(f"OCR output saved to: {ocr_output_path}")
                
                return cleaned_text
            
            return ocr_text
                
        except Exception as e:
            logger.error(f"OCR text extraction failed: {e}")
            print(f"💥 OCR text extraction failed: {e}")
            return ""
    
    def _preprocess_image_for_ocr(self, image):
        """Preprocess image to improve OCR quality"""
        try:
            from PIL import Image, ImageEnhance, ImageFilter
            
            # Convert to grayscale if not already
            if image.mode != 'L':
                image = image.convert('L')
            
            # Enhance contrast
            enhancer = ImageEnhance.Contrast(image)
            image = enhancer.enhance(1.2)
            
            # Enhance sharpness
            enhancer = ImageEnhance.Sharpness(image)
            image = enhancer.enhance(1.1)
            
            # Apply slight blur to reduce noise
            image = image.filter(ImageFilter.GaussianBlur(radius=0.5))
            
            return image
        except Exception as e:
            logger.debug(f"Image preprocessing failed: {e}")
            return image
    
    def _consolidate_ocr_results(self, results):
        """Consolidate EasyOCR results into readable text with better formatting"""
        try:
            if not results:
                return ""
            
            # Handle different result formats from EasyOCR
            processed_results = []
            for result in results:
                try:
                    # EasyOCR can return different formats: (bbox, text, confidence) or (bbox, text)
                    if len(result) == 3:
                        bbox, text, confidence = result
                    elif len(result) == 2:
                        bbox, text = result
                        confidence = 1.0  # Default confidence
                    else:
                        logger.debug(f"Unexpected result format: {result}")
                        continue
                    
                    processed_results.append((bbox, text, confidence))
                except Exception as e:
                    logger.debug(f"Error processing result: {result}, error: {e}")
                    continue
            
            if not processed_results:
                return ""
            
            # Sort results by vertical position (top to bottom)
            processed_results = sorted(processed_results, key=lambda x: x[0][0][1])  # Sort by top-left y coordinate
            
            lines = []
            current_line = []
            current_y = None
            y_tolerance = 10  # Pixels tolerance for same line
            
            for bbox, text, confidence in processed_results:
                # Skip low confidence results
                if confidence < 0.6:
                    continue
                
                # Get vertical position
                y_pos = bbox[0][1]  # top-left y coordinate
                
                # Check if this text belongs to the same line
                if current_y is None or abs(y_pos - current_y) <= y_tolerance:
                    current_line.append((bbox[0][0], text))  # (x_pos, text)
                    current_y = y_pos if current_y is None else current_y
                else:
                    # Sort current line by x position and add to lines
                    if current_line:
                        current_line.sort(key=lambda x: x[0])  # Sort by x position
                        line_text = ' '.join([text for _, text in current_line])
                        lines.append(line_text)
                    
                    # Start new line
                    current_line = [(bbox[0][0], text)]
                    current_y = y_pos
            
            # Add last line
            if current_line:
                current_line.sort(key=lambda x: x[0])
                line_text = ' '.join([text for _, text in current_line])
                lines.append(line_text)
            
            return '\n'.join(lines)
        except Exception as e:
            logger.debug(f"OCR consolidation failed: {e}")
            # Fallback to simple concatenation
            try:
                return '\n'.join([str(result[1]) if len(result) >= 2 else str(result) for result in results])
            except:
                return ""
    
    def _clean_ocr_text_enhanced(self, ocr_text: str) -> str:
        """Enhanced OCR text cleaning with better noise reduction and text consolidation"""
        import re
        
        lines = ocr_text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            line = line.strip()
            
            # Skip empty lines
            if not line:
                continue
            
            # Skip page markers
            if line.startswith('--- Page') and line.endswith('---'):
                cleaned_lines.append(line)  # Keep page markers
                continue
                
            # Skip very short lines (likely OCR noise)
            if len(line) < 2:
                continue
            
            # Skip lines that are mostly special characters
            alpha_chars = sum(1 for c in line if c.isalpha())
            if len(line) > 2 and alpha_chars < len(line) * 0.25:  # Less than 25% alphabetic characters
                continue
                
            # Skip common OCR artifacts and noise
            noise_patterns = [
                r'^[^\w\s]+$',  # Only special characters
                r'^[0-9\s\-_|]+$',  # Only numbers, spaces, dashes, underscores, pipes
                r'^[.,:;!?]+$',  # Only punctuation
                r'^www\.',  # URLs
                r'^http[s]?://',  # URLs
                r'^\d{1,3}$',  # Single numbers (page numbers, etc.)
            ]
            
            is_noise = False
            for pattern in noise_patterns:
                if re.match(pattern, line):
                    is_noise = True
                    break
            
            if is_noise:
                continue
                
            # Clean up common OCR issues
            line = re.sub(r'\s+', ' ', line)  # Multiple spaces to single space
            line = re.sub(r'[|]{2,}', ' ', line)  # Multiple pipes to space
            line = re.sub(r'[_]{3,}', ' ', line)  # Multiple underscores to space
            line = re.sub(r'[-]{3,}', ' ', line)  # Multiple dashes to space
            
            # Fix common OCR character errors
            ocr_fixes = {
                'l': 'I',  # Common l/I confusion in context
                '0': 'O',  # Common 0/O confusion in context
                '5': 'S',  # Common 5/S confusion in context
            }
            
            # Apply fixes only in specific contexts to avoid over-correction
            # This is a simplified approach - more sophisticated NLP could be used
            
            line = line.strip()
            
            # Only keep lines with meaningful content
            if len(line) >= 3 and alpha_chars >= 2:
                cleaned_lines.append(line)
        
        # Post-process to merge related lines
        merged_lines = []
        i = 0
        while i < len(cleaned_lines):
            current_line = cleaned_lines[i]
            
            # Check if current line should be merged with next line
            if (i + 1 < len(cleaned_lines) and 
                not current_line.startswith('--- Page') and
                not cleaned_lines[i + 1].startswith('--- Page') and
                len(current_line) < 80 and  # Short line
                not current_line.endswith('.') and  # Doesn't end with period
                not current_line.endswith('!') and  # Doesn't end with exclamation
                not current_line.endswith('?')):  # Doesn't end with question mark
                
                # Merge with next line
                merged_line = current_line + ' ' + cleaned_lines[i + 1]
                merged_lines.append(merged_line)
                i += 2  # Skip next line since we merged it
            else:
                merged_lines.append(current_line)
                i += 1
        
        return '\n'.join(merged_lines)
    
    def _clean_ocr_text(self, ocr_text: str) -> str:
        """Clean and filter OCR text to reduce noise and improve processing speed"""
        lines = ocr_text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            line = line.strip()
            
            # Skip empty lines
            if not line:
                continue
                
            # Skip very short lines (likely OCR noise)
            if len(line) < 3:
                continue
                
            # Skip lines that are mostly special characters or numbers only
            alpha_chars = sum(1 for c in line if c.isalpha())
            if alpha_chars < len(line) * 0.3:  # Less than 30% alphabetic characters
                continue
                
            # Skip common OCR artifacts
            if line.lower() in ['page', 'www', 'http', 'com', '|', '---', '___']:
                continue
                
            # Clean up common OCR issues
            line = line.replace('|', ' ')  # Often OCR artifacts
            line = line.replace('  ', ' ')  # Multiple spaces
            line = line.strip()
            
            if len(line) >= 5:  # Only keep lines with meaningful content
                cleaned_lines.append(line)
        
        cleaned_text = '\n'.join(cleaned_lines)
        return cleaned_text
    
    def to_dataframe(self, parsed_data: Dict[str, Any]) -> pd.DataFrame:
        return parsed_data["data"]


class WordParser(DocumentParser):
    """Parser for Word documents (DOC, DOCX)"""
    
    def parse(self, file_path: Path) -> Dict[str, Any]:
        """Parse Word document"""
        try:
            doc = Document(file_path)
            
            # Extract paragraphs
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            
            # Extract tables
            tables_data = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                if table_data:
                    # Convert to DataFrame
                    df = pd.DataFrame(table_data[1:], columns=table_data[0])
                    tables_data.append(df)
            
            # Combine data
            if tables_data:
                main_df = pd.concat(tables_data, ignore_index=True)
                data_type = "tabular"
            else:
                # Convert paragraphs to DataFrame
                main_df = pd.DataFrame({'content': paragraphs})
                data_type = "text"
            
            return {
                "type": data_type,
                "data": main_df,
                "paragraphs": paragraphs,
                "tables_found": len(tables_data),
                "metadata": {
                    "paragraph_count": len(paragraphs),
                    "table_count": len(tables_data),
                    "rows": len(main_df),
                    "columns": len(main_df.columns)
                }
            }
        except Exception as e:
            raise DocumentParseError(f"Error parsing Word document: {e}")
    
    def to_dataframe(self, parsed_data: Dict[str, Any]) -> pd.DataFrame:
        return parsed_data["data"]


class PowerPointParser(DocumentParser):
    """Parser for PowerPoint presentations (PPT, PPTX)"""
    
    def parse(self, file_path: Path) -> Dict[str, Any]:
        """Parse PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            
            slides_content = []
            tables_data = []
            
            for i, slide in enumerate(prs.slides):
                slide_content = {
                    "slide_number": i + 1,
                    "title": "",
                    "content": []
                }
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape.text_frame and shape.text_frame.paragraphs:
                            # Likely title if it's the first text or short
                            if not slide_content["title"] and len(shape.text.strip()) < 100:
                                slide_content["title"] = shape.text.strip()
                            else:
                                slide_content["content"].append(shape.text.strip())
                    
                    # Extract tables
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        
                        if table_data:
                            df = pd.DataFrame(table_data[1:], columns=table_data[0])
                            tables_data.append(df)
                
                slides_content.append(slide_content)
            
            # Create main DataFrame
            if tables_data:
                main_df = pd.concat(tables_data, ignore_index=True)
                data_type = "tabular"
            else:
                # Convert slides to text format
                slide_rows = []
                for slide in slides_content:
                    content_text = " ".join(slide["content"])
                    slide_rows.append({
                        "slide": slide["slide_number"],
                        "title": slide["title"],
                        "content": content_text
                    })
                main_df = pd.DataFrame(slide_rows)
                data_type = "presentation"
            
            return {
                "type": data_type,
                "data": main_df,
                "slides": slides_content,
                "tables_found": len(tables_data),
                "metadata": {
                    "slide_count": len(slides_content),
                    "table_count": len(tables_data),
                    "rows": len(main_df),
                    "columns": len(main_df.columns)
                }
            }
        except Exception as e:
            raise DocumentParseError(f"Error parsing PowerPoint: {e}")
    
    def to_dataframe(self, parsed_data: Dict[str, Any]) -> pd.DataFrame:
        return parsed_data["data"]


class TextParser(DocumentParser):
    """Parser for plain text files (TXT)"""
    
    def parse(self, file_path: Path) -> Dict[str, Any]:
        """Parse text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Split into paragraphs/lines for processing
            lines = [line.strip() for line in content.split('\n') if line.strip()]
            
            # Create DataFrame with content
            main_df = pd.DataFrame({'content': lines})
            
            return {
                "type": "text",
                "data": main_df,
                "raw_text": content,
                "metadata": {
                    "line_count": len(lines),
                    "character_count": len(content),
                    "rows": len(main_df),
                    "columns": len(main_df.columns)
                }
            }
        except Exception as e:
            raise DocumentParseError(f"Error parsing text file: {e}")
    
    def to_dataframe(self, parsed_data: Dict[str, Any]) -> pd.DataFrame:
        return parsed_data["data"]


class DocumentParserFactory:
    """Factory for creating appropriate document parsers"""
    
    PARSER_MAP = {
        FileFormat.CSV: CSVParser,
        FileFormat.XLSX: ExcelParser,
        FileFormat.XLS: ExcelParser,
        FileFormat.PDF: PDFParser,
        FileFormat.DOC: WordParser,
        FileFormat.DOCX: WordParser,
        FileFormat.PPT: PowerPointParser,
        FileFormat.PPTX: PowerPointParser,
        FileFormat.TXT: TextParser,
    }
    
    @classmethod
    def create_parser(cls, file_path: Path) -> DocumentParser:
        """Create appropriate parser for file format"""
        file_format = FileDetector.detect_format(file_path)
        
        if file_format not in cls.PARSER_MAP:
            raise DocumentParseError(f"Unsupported file format: {file_format.value}")
        
        parser_class = cls.PARSER_MAP[file_format]
        return parser_class()
    
    @classmethod
    def parse_document(cls, file_path: Path) -> Dict[str, Any]:
        """Parse any supported document format"""
        parser = cls.create_parser(file_path)
        return parser.parse(file_path)
    
    @classmethod
    def document_to_dataframe(cls, file_path: Path) -> pd.DataFrame:
        """Parse document and convert to DataFrame"""
        parser = cls.create_parser(file_path)
        parsed_data = parser.parse(file_path)
        return parser.to_dataframe(parsed_data)